P&G ESG Intelligence Platform — AI Kit A Demo (v3)
- Executive Dashboard with Supplier Portal status
- CSRD Analyzer (supplier linkage, expired cert penalty)
- Narrative Generator with optional LLM + trace
- Waste Optimizer with Anomaly Detection (IsolationForest) and ML-backed Savings Simulator
- Supplier Data Portal with evidence uploads + magic links
- Assurance Readiness checklist
- Board Presentation PPTX export

Note: This is a consolidated demo app you can run as-is.
"""

import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
from datetime import datetime, timedelta
from io import BytesIO
import base64
import time

# Optional: PPTX export
try:
    from pptx import Presentation
except Exception:
    Presentation = None

# Optional ML libs (demo-safe)
try:
    from sklearn.ensemble import IsolationForest, GradientBoostingRegressor
    from sklearn.model_selection import train_test_split
    from sklearn.metrics import r2_score
except Exception:
    IsolationForest = None
    GradientBoostingRegressor = None
    train_test_split = None
    r2_score = None

st.set_page_config(page_title="P&G ESG Intelligence Platform", page_icon="🌍", layout="wide")

# -------------------------
# Session state & branding
# -------------------------
ss = st.session_state
ss.setdefault("selected_module", "Executive Dashboard")
ss.setdefault("csrd_data", None)
ss.setdefault("waste_data", None)
ss.setdefault("supplier_db", {"suppliers": [], "materials": [], "svhc": [], "circularity": [], "evidence": [], "magic_links": []})

BRANDS = {
    "P&G": {"primary": "#003DA5", "logo": "https://upload.wikimedia.org/wikipedia/commons/thumb/f/f9/Procter_%26_Gamble_logo.svg/240px-Procter_%26_Gamble_logo.svg.png"},
    "Unilever": {"primary": "#1F45FC", "logo": "https://upload.wikimedia.org/wikipedia/commons/thumb/4/4f/Unilever.svg/240px-Unilever.svg.png"},
    "L'Oréal": {"primary": "#000000", "logo": "https://upload.wikimedia.org/wikipedia/commons/thumb/a/ab/L%27Oreal_logo.svg/320px-L%27Oreal_logo.svg.png"},
}
ss.setdefault("branding", {"client":"P&G", **BRANDS["P&G"]})

# -------------------------
# Utility helpers
# -------------------------

def calc_completeness(df: pd.DataFrame, topic: str) -> float:
    if df is None or df.empty:
        return 0.0
    return round((df.notna().sum().sum() / (len(df)*len(df.columns))) * 100, 1)


def supplier_coverage():
    db = ss.supplier_db
    mats_total = len(db["materials"]) or 0
    e2_cov = len({r["material_id"] for r in db["svhc"]})/mats_total*100 if mats_total else 0.0
    e5_cov = len({r["material_id"] for r in db["circularity"]})/mats_total*100 if mats_total else 0.0
    return round(e2_cov,1), round(e5_cov,1)


def expired_certificates():
    out = []
    today = datetime.now().date()
    for row in ss.supplier_db.get("circularity", []):
        try:
            exp = pd.to_datetime(str(row.get("cert_expiry")), errors="coerce").date()
            if exp <= today:
                out.append(row)
        except Exception:
            out.append(row)
    return out


def save_evidence(uploaded_file):
    if uploaded_file is None:
        return None
    eid = f"EVD-{len(ss.supplier_db['evidence'])+1:05d}"
    ss.supplier_db["evidence"].append({
        "id": eid,
        "name": uploaded_file.name,
        "size": getattr(uploaded_file, 'size', None),
        "uploaded_at": datetime.now().isoformat(timespec='seconds'),
        "bytes": uploaded_file.getbuffer().tobytes(),
        "mime": "application/pdf",
    })
    return eid

# Feature engineering for ML

def _featureize_waste(df: pd.DataFrame) -> pd.DataFrame:
    if df is None or df.empty:
        return pd.DataFrame()
    x = df.copy()
    if 'cost_eur' not in x.columns and 'cost_usd' in x.columns:
        x['cost_eur'] = pd.to_numeric(x['cost_usd'], errors='coerce')
    for c in ['waste_kg','production_volume','cost_eur']:
        if c in x.columns:
            x[c] = pd.to_numeric(x[c], errors='coerce')
    if 'date' in x.columns:
        dt = pd.to_datetime(x['date'], errors='coerce')
        x['month'] = dt.dt.month.fillna(0).astype(int)
        x['dow'] = dt.dt.dayofweek.fillna(0).astype(int)
    else:
        x['month'] = 0; x['dow'] = 0
    for col in ['product_line','plant','waste_type','waste_reason','disposal_method']:
        if col in x.columns:
            d = pd.get_dummies(x[col], prefix=col, dummy_na=False)
            keep = d.sum().sort_values(ascending=False).head(5).index
            x = pd.concat([x, d[keep]], axis=1)
    x['waste_rate'] = (x.get('waste_kg',0) / x.get('production_volume',1).replace(0,np.nan)).fillna(0) * 100
    return x.fillna(0)

# Demo data generators

def generate_sample_waste_data():
    np.random.seed(42)
    dates = pd.date_range('2024-01-01', '2024-12-31', freq='D')
    products = ['Head & Shoulders', 'Pantene', 'Olay', 'SK-II', 'Gillette']
    plants = ['Plant A - France', 'Plant B - Germany', 'Plant C - Belgium']
    recs = []
    for d in dates:
        for _ in range(np.random.randint(5, 10)):
            recs.append({
                'date': d,
                'product_line': np.random.choice(products),
                'plant': np.random.choice(plants),
                'batch_id': f"BATCH-{np.random.randint(1000,9999)}",
                'production_volume': np.random.uniform(1000, 10000),
                'waste_kg': np.random.uniform(50, 500),
                'waste_type': np.random.choice(['Packaging','Raw Material','Product Defects','Cleaning']),
                'waste_reason': np.random.choice(['Quality Issue','Changeover','Expiry','Equipment Failure','Human Error']),
                'disposal_method': np.random.choice(['Recycled','Incinerated','Landfill','Composted']),
                'cost_eur': np.random.uniform(100, 2000),
            })
    return pd.DataFrame(recs)


def generate_sample_csrd_data():
    np.random.seed(42)
    dates = pd.date_range('2021-01-01', '2023-12-31', freq='M')
    data = {
        'date': dates,
        'facility': np.random.choice(['Plant A - France','Plant B - Germany','Plant C - Belgium'], len(dates)),
        'total_water_withdrawal': np.random.uniform(1000, 5000, len(dates)),
        'water_consumption': np.random.uniform(800, 4000, len(dates)),
        'water_discharge': np.random.uniform(200, 1000, len(dates)),
        'water_recycled_percentage': np.random.uniform(20, 60, len(dates)),
        'plastic_packaging_total': np.random.uniform(100, 500, len(dates)),
        'recycled_content_percentage': np.random.uniform(10, 40, len(dates)),
        'recyclable_packaging_percentage': np.random.uniform(60, 95, len(dates)),
        'hazardous_waste': np.random.uniform(10, 100, len(dates)),
        'non_hazardous_waste': np.random.uniform(200, 1000, len(dates)),
        'waste_to_landfill': np.random.uniform(50, 300, len(dates)),
        'waste_recycled': np.random.uniform(100, 600, len(dates)),
        'air_emissions_nox': np.random.uniform(1, 20, len(dates)),
        'air_emissions_sox': np.random.uniform(0.5, 10, len(dates)),
        'air_emissions_pm': np.random.uniform(0.1, 5, len(dates)),
        'product_carbon_footprint': np.random.uniform(0.5, 2.5, len(dates)),
    }
    return pd.DataFrame(data)

# -------------------------
# UI: Sidebar
# -------------------------
with st.sidebar:
    # Branding/theme
    brand = st.selectbox("Client theme", list(BRANDS.keys()), index=[*BRANDS].index(ss.branding["client"]))
    ss.branding = {"client": brand, **BRANDS[brand]}
    st.image(ss.branding["logo"], width=140)
    st.markdown(f"""
    <style>
      .stButton>button {{ background-color: {ss.branding['primary']} !important; color: white; font-weight: 600; }}
      .header-style {{ background: linear-gradient(90deg, {ss.branding['primary']} 0%, #0052CC 100%) !important; }}
    </style>
    """, unsafe_allow_html=True)

    st.markdown("---")
    st.markdown("### 📊 Navigation")
    module = st.radio("", [
        "Executive Dashboard",
        "CSRD Compliance Analyzer",
        "Waste Optimizer",
        "Supplier Data Portal",
        "Assurance Readiness",
        "Report Generator",
    ])
    ss.selected_module = module

    st.markdown("---")
    if st.button("Load sample supplier data"):
        if not ss.supplier_db["suppliers"]:
            # seed minimal sample set
            sid = "SUP-00001"; mid1 = "MAT-00001"; mid2 = "MAT-00002"
            ss.supplier_db["suppliers"].append({"id":sid, "name":"GreenPoly GmbH","country":"DE","email":"info@greenpoly.de"})
            ss.supplier_db["materials"].extend([
                {"id": mid1, "supplier_id": sid, "name":"PET Bottle Resin","cas_ec":"25038-59-9","sku":"PG-BTL-500","unit_mass_g":30.0},
                {"id": mid2, "supplier_id": sid, "name":"HDPE Cap","cas_ec":"25213-02-9","sku":"PG-CAP-28HD","unit_mass_g":5.5},
            ])
            today = datetime.now().date()
            ss.supplier_db["svhc"].append({"material_id": mid1, "present": False, "svhc_list": [], "max_pct":0.0, "scip":"SCIP-123","reach_ref":"REACH-REF-001", "declaration_date": str(today), "expiry_date": str(today.replace(year=today.year+1)), "evidence": None})
            ss.supplier_db["circularity"].append({"material_id": mid1, "pcr_pct": 30.0, "pre_consumer_pct":10.0, "method":"physical", "cert_type":"UL 2809","cert_id":"UL-12345","issuer":"UL","cert_expiry": str(today.replace(year=today.year+1)), "evidence": None})
            ss.supplier_db["svhc"].append({"material_id": mid2, "present": True, "svhc_list":["D4"], "max_pct":0.12, "scip":"SCIP-456","reach_ref":"REACH-REF-002", "declaration_date": str(today), "expiry_date": str(today.replace(year=today.year+1)), "evidence": None})
            ss.supplier_db["circularity"].append({"material_id": mid2, "pcr_pct": 20.0, "pre_consumer_pct":5.0, "method":"mass_balance", "cert_type":"RecyClass","cert_id":"RC-9876","issuer":"RecyClass","cert_expiry": str(today.replace(year=today.year-1)), "evidence": None})
            st.success("Sample supplier data loaded.")
        else:
            st.info("Supplier/material records already present — skipped.")

# -------------------------
# Header
# -------------------------
st.markdown(f"""
<div class="header-style" style='padding:1.5rem;border-radius:10px;color:white;margin-bottom:1rem;'>
  <h1 style='margin:0'>🌍 P&G ESG Intelligence Platform</h1>
  <p style='margin:6px 0 0 0; font-size:18px'>AI-Powered CSRD Compliance & Waste Optimization</p>
</div>
""", unsafe_allow_html=True)

# -------------------------
# Executive Dashboard
# -------------------------
if ss.selected_module == "Executive Dashboard":
    c1,c2,c3,c4 = st.columns(4)
    c1.metric("Fine Risk Identified", "€1.8M", "+5%")
    c2.metric("CSRD Compliance", "78%", "+5%")
    c3.metric("Waste Savings Potential", "€2.3M", "+12%")
    c4.metric("Waste Diverted", "92%", "✔")

    st.markdown("---")
    st.subheader("🔗 Supplier Portal Status")
    sp1, sp2, sp3, sp4, sp5, sp6 = st.columns(6)
    total_sup = len(ss.supplier_db.get("suppliers", []))
    total_mat = len(ss.supplier_db.get("materials", []))
    e2_cov, e5_cov = supplier_coverage()
    links = len(ss.supplier_db.get("magic_links", []))
    exp_certs = len(expired_certificates())
    sp1.metric("Suppliers", total_sup)
    sp2.metric("Materials", total_mat)
    sp3.metric("SVHC coverage", f"{e2_cov}%")
    sp4.metric("PCR coverage", f"{e5_cov}%")
    sp5.metric("Magic links", links)
    sp6.metric("Expired certs", exp_certs)

# -------------------------
# CSRD Compliance Analyzer
# -------------------------
elif ss.selected_module == "CSRD Compliance Analyzer":
    tab1, tab2, tab3 = st.tabs(["📤 Data Upload", "🎯 Gap Analysis", "📝 Narrative Generator"])

    with tab1:
        st.markdown("### Upload Your ESG Data")
        up = st.file_uploader("CSV or Excel", type=["csv","xlsx"])
        if st.button("🔬 Use Demo Data"):
            ss.csrd_data = generate_sample_csrd_data()
            st.success("Demo data loaded.")
        if up is not None:
            ss.csrd_data = pd.read_csv(up) if up.name.endswith('.csv') else pd.read_excel(up)
            st.success("Data uploaded.")
        if ss.csrd_data is not None:
            st.dataframe(ss.csrd_data.head(), use_container_width=True)
            st.info(f"Completeness: {calc_completeness(ss.csrd_data, 'All')}%")

    with tab2:
        st.markdown("### 🎯 ESRS Gap Analysis")
        if ss.csrd_data is None:
            st.warning("Please upload or load demo data.")
        else:
            topic = st.selectbox("ESRS Topic", ["Water","Circular Economy","Pollution"]) 
            if st.button("🔍 Analyze Gaps", type="primary"):
                gap_catalog = {
                    "Water": [
                        {"requirement":"ESRS E3-1: Water management policies","status":"⚠️ Partial","gap":"Missing water stress assessment for 3 facilities","risk":"Medium","fine_max":200_000},
                        {"requirement":"ESRS E3-4: Water consumption intensity","status":"❌ Missing","gap":"No product-level water intensity metrics","risk":"High","fine_max":500_000},
                        {"requirement":"ESRS E3-5: Water discharge quality","status":"✅ Complete","gap":"None","risk":"Low","fine_max":0},
                    ],
                    "Circular Economy": [
                        {"requirement":"ESRS E5-1: Resource inflows","status":"⚠️ Partial","gap":"Missing supplier recycled content data for 40% of materials","risk":"High","fine_max":700_000},
                        {"requirement":"ESRS E5-2: Resource outflows","status":"✅ Complete","gap":"None","risk":"Low","fine_max":0},
                        {"requirement":"ESRS E5-5: Circular design","status":"❌ Missing","gap":"No systematic circularity assessment for new products","risk":"Medium","fine_max":300_000},
                    ],
                    "Pollution": [
                        {"requirement":"ESRS E2-1: Pollution policies","status":"✅ Complete","gap":"None","risk":"Low","fine_max":0},
                        {"requirement":"ESRS E2-4: Pollution of air","status":"⚠️ Partial","gap":"Missing Scope 3 air emission data from logistics","risk":"Medium","fine_max":400_000},
                        {"requirement":"ESRS E2-6: Substances of concern","status":"❌ Missing","gap":"Incomplete SVHC tracking in supply chain","risk":"High","fine_max":1_000_000},
                    ],
                }
                gaps = gap_catalog.get(topic, [])
                total_fine = 0; high = 0
                for g in gaps:
                    em = "🔴" if g['risk']=="High" else "🟡" if g['risk']=="Medium" else "🟢"
                    with st.expander(f"{em} {g['requirement']} - {g['status']}"):
                        a,b = st.columns([2,1])
                        a.markdown(f"**Gap:** {g['gap']}")
                        a.markdown(f"**Risk:** {g['risk']}")
                        b.markdown(f"**Fine Risk (max):** €{g['fine_max']:,}")
                    total_fine += g['fine_max']
                    if g['risk']=="High": high += 1
                e2_cov, e5_cov = supplier_coverage()
                expired = expired_certificates()
                expired_pen = min(len(expired)*2, 10)
                base_score = max(0, 100 - high*15)
                bonus = round((e2_cov + e5_cov) / 20.0, 1)  # up to +10
                final = max(0, min(100, base_score + bonus - expired_pen))
                st.markdown("---")
                c1,c2,c3,c4 = st.columns(4)
                c1.metric("High Risks", high)
                c2.metric("Base Score", f"{base_score}%")
                c3.metric("Coverage Bonus", f"+{bonus} pts")
                c4.metric("Expired Certs Penalty", f"-{expired_pen} pts")
                st.success(f"Adjusted Compliance Score: {final}%")
                if expired:
                    st.markdown("#### Certificates needing renewal")
                    exp_df = pd.DataFrame(expired)
                    st.dataframe(exp_df[["material_id","cert_type","cert_id","issuer","cert_expiry"]], use_container_width=True)

    with tab3:
        st.markdown("### 📝 Disclosure Narrative Generator (Demo + Optional LLM)")
        if ss.csrd_data is None:
            st.warning("Please upload data or load demo data.")
        else:
            topic = st.selectbox("Topic", ["Water","Circular Economy","Pollution"], key="nar_topic")
            audit_ready = st.checkbox("This data element is audit-ready", True)
            word_count = st.select_slider("Word Count", options=[300,500,1000,1500], value=500)
            st.markdown("#### 🔎 Trace inputs that will feed the narrative")
            trace = {"records": len(ss.csrd_data), "columns": list(ss.csrd_data.columns)[:10], "completeness": calc_completeness(ss.csrd_data, topic)}
            st.json(trace)
            with st.expander("⚙️ Optional: Use LLM (Anthropic/OpenAI)"):
                provider = st.selectbox("Provider", ["None","Anthropic","OpenAI"], index=0)
                api_key = st.text_input("API Key", type="password")
                model = st.text_input("Model (e.g., claude-3-5-sonnet or gpt-4o-mini)")
                temperature = st.slider("Creativity", 0.0, 1.0, 0.2)
            if st.button("📝 Generate Narrative", type="primary"):
                if not audit_ready:
                    st.error("❌ Cannot generate narrative - mark audit-ready first.")
                else:
                    prompt = f"Generate a {word_count}-word ESRS narrative for {topic}. Completeness: {trace['completeness']}%. Demo values; keep audit-conscious."
                    content = None
                    if provider != "None" and api_key:
                        try:
                            if provider=="Anthropic":
                                from anthropic import Anthropic
                                client = Anthropic(api_key=api_key)
                                msg = client.messages.create(model=(model or "claude-3-5-sonnet-20240620"), max_tokens=700, temperature=temperature, messages=[{"role":"user","content":prompt}])
                                content = "\n\n".join([getattr(b, "text", str(b)) for b in msg.content])
                            elif provider=="OpenAI":
                                from openai import OpenAI
                                oai = OpenAI(api_key=api_key)
                                chat = oai.chat.completions.create(model=(model or "gpt-4o-mini"), messages=[{"role":"user","content":prompt}], temperature=temperature, max_tokens=700)
                                content = chat.choices[0].message.content
                        except Exception as e:
                            st.warning(f"LLM call failed; falling back to demo text. ({e})")
                    if not content:
                        content = f"""
**{topic} (Illustrative ESRS Narrative)**

Governance & Strategy: oversight via ESG Committee; alignment with SDGs.

Risk & Opportunity: facility-level assessments; investment roadmap for recycling and efficiency.

Metrics & Targets: demo-only values. Replace with assured data before disclosure.

*Word target:* ~{word_count} | *Completeness:* {trace['completeness']}%
"""
                    st.markdown("### Generated ESRS Narrative")
                    st.markdown(content)
                    st.download_button("📥 Download Narrative (Markdown)", content, file_name=f"ESRS_{topic}_Disclosure_DEMO.md", mime="text/markdown")

# -------------------------
# Waste Optimizer
# -------------------------
elif ss.selected_module == "Waste Optimizer":
    tab1, tab2, tab3, tab4 = st.tabs(["📤 Data Input", "📊 Analysis", "💰 Savings Simulator", "📈 Insights"])    
    with tab1:
        up = st.file_uploader("Upload waste data (CSV/Excel)", type=["csv","xlsx"])
        if st.button("🔬 Generate Demo Data"):
            ss.waste_data = generate_sample_waste_data()
            st.success("Demo waste data generated.")
        if up is not None:
            ss.waste_data = pd.read_csv(up) if up.name.endswith('.csv') else pd.read_excel(up)
            st.success("Data uploaded.")
        if ss.waste_data is not None:
            st.dataframe(ss.waste_data.head(), use_container_width=True)

    with tab2:
        st.markdown("### 📊 Analysis")
        if ss.waste_data is None:
            st.warning("Upload or generate demo data.")
        else:
            # Waste by product and cause
            by_product = ss.waste_data.groupby('product_line')['waste_kg'].sum().sort_values(ascending=False)
            by_reason = ss.waste_data.groupby('waste_reason')['waste_kg'].sum().sort_values(ascending=False)
            a,b = st.columns(2)
            with a:
                st.subheader("Waste by Product Line")
                fig = px.bar(x=by_product.head(5).values, y=by_product.head(5).index, orientation='h')
                fig.update_layout(height=300, showlegend=False)
                st.plotly_chart(fig, use_container_width=True)
            with b:
                st.subheader("Waste by Root Cause")
                fig = px.pie(values=by_reason.values, names=by_reason.index)
                fig.update_layout(height=300)
                st.plotly_chart(fig, use_container_width=True)

            st.subheader("Waste Heatmap: Product × Plant")
            pivot = ss.waste_data.pivot_table(values='waste_kg', index='product_line', columns='plant', aggfunc='sum')
            fig = px.imshow(pivot, aspect='auto', color_continuous_scale='Reds')
            fig.update_layout(height=400)
            st.plotly_chart(fig, use_container_width=True)

            # Anomaly Detection
            st.markdown("### 🔎 Anomaly Detection (IsolationForest)")
            if IsolationForest is None:
                st.info("Install scikit-learn for full anomaly detection (`pip install scikit-learn`). Showing heuristic spike finder.")
                _df = ss.waste_data.copy(); _df['date'] = pd.to_datetime(_df['date'], errors='coerce')
                daily = _df.groupby('date')['waste_kg'].sum().reset_index()
                thr = daily['waste_kg'].median() + 2.5 * daily['waste_kg'].mad()
                daily['anomaly'] = daily['waste_kg'] > thr
            else:
                _X = _featureize_waste(ss.waste_data)
                _X['date'] = pd.to_datetime(ss.waste_data['date'], errors='coerce')
                daily = _X.groupby('date').agg({'waste_kg':'sum','waste_rate':'mean','month':'first','dow':'first'}).reset_index()
                model = IsolationForest(n_estimators=200, contamination=0.05, random_state=42)
                Z = daily[['waste_kg','waste_rate','month','dow']].values
                daily['score'] = model.fit_predict(Z)
                daily['anomaly'] = daily['score'] == -1
            fig_an = go.Figure()
            fig_an.add_trace(go.Scatter(x=daily['date'], y=daily['waste_kg'], name='Daily waste', mode='lines+markers'))
            if daily['anomaly'].any():
                spikes = daily[daily['anomaly']]
                fig_an.add_trace(go.Scatter(x=spikes['date'], y=spikes['waste_kg'], name='Anomalies', mode='markers', marker=dict(size=10, symbol='x')))
            fig_an.update_layout(height=320, title='Daily Waste with Anomalies')
            st.plotly_chart(fig_an, use_container_width=True)
            try:
                corr = daily.assign(anom=daily['anomaly'].astype(int))[['anom','waste_kg','waste_rate','month','dow']].corr()['anom'].drop('anom').abs().sort_values(ascending=False)
                st.dataframe(corr.to_frame('|corr|'))
            except Exception:
                pass

    with tab3:
        st.markdown("### 💰 Savings Simulator")
        if ss.waste_data is None:
            st.warning("Upload waste data first.")
        else:
            l, r = st.columns(2)
            with l:
                st.markdown("**Process Improvements**")
                process_reduction = st.slider("Process optimization (%)", 0, 50, 15)
                equipment_upgrade = st.checkbox("Equipment upgrade (−20% changeover waste)", True)
                predictive_maintenance = st.checkbox("Predictive maintenance (−30% equipment failures)", True)
                st.markdown("**Material Changes**")
                material_substitution = st.slider("Material substitution savings (%)", 0, 30, 10)
                packaging_redesign = st.checkbox("Packaging redesign (−25% packaging waste)", True)
            with r:
                st.markdown("**Specification Changes**")
                spec_tolerance = st.slider("Specification tolerance (%)", 0, 20, 8)
                quality_improvement = st.checkbox("Quality system upgrade (−40% defects)", True)
                st.markdown("**Operational Changes**")
                training_program = st.checkbox("Enhanced operator training (−15% human error)", True)
                inventory_optimization = st.checkbox("Inventory optimization (−50% expiry waste)", True)
                use_ml = st.checkbox("Use ML-backed estimate (if sklearn present)", True)

            if st.button("💡 Calculate Savings", type="primary"):
                df = ss.waste_data
                base_cost = float(pd.to_numeric(df.get('cost_eur', pd.Series(dtype=float)), errors='coerce').fillna(0).sum())
                savings = 0.0
                savings += base_cost * (process_reduction / 100)
                if equipment_upgrade: savings += base_cost * 0.10
                if predictive_maintenance: savings += base_cost * 0.08
                savings += base_cost * (material_substitution / 100)
                if packaging_redesign: savings += base_cost * 0.12
                savings += base_cost * (spec_tolerance / 100)
                if quality_improvement: savings += base_cost * 0.15
                if training_program: savings += base_cost * 0.05
                if inventory_optimization: savings += base_cost * 0.03

                # ML uplift
                if use_ml and GradientBoostingRegressor is not None:
                    try:
                        X = _featureize_waste(df)
                        y = pd.to_numeric(df.get('cost_eur', pd.Series(dtype=float)), errors='coerce').fillna(0)
                        common = X.index.intersection(y.index)
                        X, y = X.loc[common], y.loc[common]
                        if len(X) > 50 and X.select_dtypes(include=[np.number]).shape[1] >= 5:
                            feats = X.select_dtypes(include=[np.number])
                            Xtr, Xte, ytr, yte = train_test_split(feats, y, test_size=0.2, random_state=42)
                            gbr = GradientBoostingRegressor(random_state=42).fit(Xtr, ytr)
                            base_pred = gbr.predict(feats).sum()
                            # simple counterfactual scaling
                            cf = feats.copy()
                            cf_factor = 1.0
                            cf_factor *= (1 - process_reduction/100)
                            if equipment_upgrade: cf_factor *= 0.90
                            if predictive_maintenance: cf_factor *= 0.92
                            if packaging_redesign: cf_factor *= 0.88
                            if quality_improvement: cf_factor *= 0.85
                            if training_program: cf_factor *= 0.95
                            if inventory_optimization: cf_factor *= 0.97
                            if spec_tolerance: cf_factor *= (1 - spec_tolerance/100)
                            for col in [c for c in cf.columns if c.startswith('waste_') or c in ['waste_kg','waste_rate']]:
                                cf[col] = cf[col] * cf_factor
                            cf_pred = gbr.predict(cf).sum()
                            ml_delta = max(0.0, base_pred - cf_pred)
                            if ml_delta > 0:
                                savings = ml_delta
                            if r2_score is not None:
                                st.caption(f"ML uplift used (GradientBoostingRegressor). Test R² ≈ {r2_score(yte, gbr.predict(Xte)):.2f} (demo).")
                    except Exception as e:
                        st.info(f"ML uplift skipped ({e}).")

                st.markdown("---")
                st.markdown("### 📊 Projected Annual Savings (Demo)")
                c1,c2,c3 = st.columns(3)
                c1.markdown(f"<div class='success-box'><h3>💰 Total Savings</h3><p style='font-size:28px;font-weight:bold'>€{savings/1_000_000:.2f}M</p><p>Per year</p></div>", unsafe_allow_html=True)
                reduction_pct = (savings / base_cost * 100) if base_cost else 0.0
                c2.markdown(f"<div class='success-box'><h3>📉 Waste Reduction</h3><p style='font-size:28px;font-weight:bold'>{reduction_pct:.1f}%</p><p>Overall reduction</p></div>", unsafe_allow_html=True)
                roi = ((savings - 500_000) / 500_000 * 100) if 500_000 else 0.0
                c3.markdown(f"<div class='success-box'><h3>📈 ROI</h3><p style='font-size:28px;font-weight:bold'>{roi:.0f}%</p><p>12-month ROI (assumes €500K investment)</p></div>", unsafe_allow_html=True)

# -------------------------
# Supplier Data Portal
# -------------------------
elif ss.selected_module == "Supplier Data Portal":
    st.markdown("## 🔗 Supplier Data Portal (Demo)")
    st.caption("Collect SVHC declarations and recycled-content certificates from vendors.")
    tabA, tabB, tabC = st.tabs(["➕ Suppliers & Materials","📄 Declarations","🔗 Magic Links"])

    with tabA:
        st.subheader("Suppliers")
        c1,c2,c3 = st.columns(3)
        with c1: s_name = st.text_input("Supplier name")
        with c2: s_country = st.text_input("Country")
        with c3: s_email = st.text_input("Contact email")
        if st.button("Add Supplier") and s_name:
            sid = f"SUP-{len(ss.supplier_db['suppliers'])+1:05d}"
            ss.supplier_db["suppliers"].append({"id": sid, "name": s_name, "country": s_country, "email": s_email})
            st.success(f"Added supplier {s_name} ({sid})")
        st.subheader("Materials")
        sup_opts = {s["name"]: s["id"] for s in ss.supplier_db["suppliers"]}
        mcol1,mcol2,mcol3,mcol4 = st.columns([2,1,1,1])
        with mcol1: sel_sup = st.selectbox("Supplier", list(sup_opts.keys()) or ["— add a supplier first —"]) 
        with mcol2: m_name = st.text_input("Material name")
        with mcol3: m_sku = st.text_input("SKU")
        with mcol4: m_casec = st.text_input("CAS/EC")
        um = st.number_input("Unit mass (g)", 0.0, 1e6, 100.0)
        if st.button("Add Material") and sup_opts:
            mid = f"MAT-{len(ss.supplier_db['materials'])+1:05d}"
            ss.supplier_db["materials"].append({"id": mid, "supplier_id": sup_opts.get(sel_sup, ""), "name": m_name, "cas_ec": m_casec, "sku": m_sku, "unit_mass_g": um})
            st.success(f"Added material {m_name} ({mid})")
        st.markdown("### Current Records")
        s_df = pd.DataFrame(ss.supplier_db["suppliers"]) if ss.supplier_db["suppliers"] else pd.DataFrame(columns=["id","name","country","email"])
        m_df = pd.DataFrame(ss.supplier_db["materials"]) if ss.supplier_db["materials"] else pd.DataFrame(columns=["id","supplier_id","name","cas_ec","sku","unit_mass_g"])
        st.dataframe(s_df, use_container_width=True)
        st.dataframe(m_df, use_container_width=True)

    with tabB:
        st.subheader("SVHC (E2) and Recycled Content (E5)")
        mats = {f"{m['id']} – {m['name']} ({m['sku']})": m['id'] for m in ss.supplier_db["materials"]}
        mat_key = st.selectbox("Material", list(mats.keys()) or ["— add a material first —"]) 
        c1, c2 = st.columns(2)
        with c1:
            st.markdown("**SVHC Declaration**")
            present = st.selectbox("SVHC present?", ["No","Yes"]) 
            svhc_text = st.text_input("SVHC list (comma-separated)")
            max_pct = st.number_input("Max concentration (% w/w)", 0.0, 100.0, 0.0)
            scip = st.text_input("SCIP ID")
            reach_ref = st.text_input("REACH reference")
            decl_date = st.date_input("Declaration date")
            expiry = st.date_input("Declaration expiry")
            svhc_file = st.file_uploader("Upload SVHC evidence (PDF)", type=["pdf"], key="svhc_file")
        with c2:
            st.markdown("**Recycled Content**")
            pcr = st.number_input("Post-consumer recycled (%)", 0.0, 100.0, 0.0)
            pre = st.number_input("Pre-consumer recycled (%)", 0.0, 100.0, 0.0)
            method = st.selectbox("Method", ["mass_balance","physical"]) 
            cert_type = st.selectbox("Certification", ["UL 2809","RecyClass","SCS","ISCC PLUS","None"]) 
            cert_id = st.text_input("Certificate ID")
            issuer = st.text_input("Issuer")
            cert_expiry = st.date_input("Certificate expiry")
            circ_file = st.file_uploader("Upload certification (PDF)", type=["pdf"], key="circ_file")
        if st.button("Save Declarations") and mats:
            svhc_eid = save_evidence(svhc_file)
            circ_eid = save_evidence(circ_file)
            ss.supplier_db["svhc"].append({
                "material_id": mats.get(mat_key, ""),
                "present": present=="Yes",
                "svhc_list": [s.strip() for s in svhc_text.split(',') if s.strip()],
                "max_pct": max_pct,
                "scip": scip,
                "reach_ref": reach_ref,
                "declaration_date": str(decl_date),
                "expiry_date": str(expiry),
                "evidence": svhc_eid,
            })
            ss.supplier_db["circularity"].append({
                "material_id": mats.get(mat_key, ""),
                "pcr_pct": pcr,
                "pre_consumer_pct": pre,
                "method": method,
                "cert_type": cert_type,
                "cert_id": cert_id,
                "issuer": issuer,
                "cert_expiry": str(cert_expiry),
                "evidence": circ_eid,
            })
            st.success("Saved.")

        st.markdown("### Evidence Uploaded (this session)")
        ev_df = pd.DataFrame(ss.supplier_db["evidence"]) if ss.supplier_db["evidence"] else pd.DataFrame(columns=["id","name","size","uploaded_at"])
        st.dataframe(ev_df[["id","name","size","uploaded_at"]], use_container_width=True)
        if not ev_df.empty:
            sel_eid = st.selectbox("Preview evidence", ev_df["id"].tolist())
            row = next((e for e in ss.supplier_db["evidence"] if e["id"] == sel_eid), None)
            if row and row.get("bytes"):
                st.download_button("⬇️ Download selected PDF", data=row["bytes"], file_name=row["name"], mime=row.get("mime","application/pdf"))
                b64 = base64.b64encode(row["bytes"]).decode()
                st.markdown(f"<iframe src='data:application/pdf;base64,{b64}' width='100%' height='500' style='border:1px solid #eee;border-radius:8px'></iframe>", unsafe_allow_html=True)

    with tabC:
        st.subheader("Magic Links (Demo)")
        st.caption("Generate shareable tokens to simulate vendor self-serve data entry.")
        sup_map = {s['name']: s['id'] for s in ss.supplier_db['suppliers']}
        sel_sup_for_link = st.selectbox("Supplier", list(sup_map.keys()) or ["— none —"]) 
        if st.button("Create supplier link") and sup_map:
            token = f"DEMOLINK-{len(ss.supplier_db['magic_links'])+1:05d}"
            ss.supplier_db['magic_links'].append({"token": token, "type": "supplier", "target_id": sup_map.get(sel_sup_for_link, "")})
            st.success(f"Created token: {token}")
        mat_map = {f"{m['name']} ({m['sku']})": m['id'] for m in ss.supplier_db['materials']}
        sel_mat_for_link = st.selectbox("Material", list(mat_map.keys()) or ["— none —"], key="mat_link")
        if st.button("Create material link") and mat_map:
            token = f"DEMOLINK-{len(ss.supplier_db['magic_links'])+1:05d}"
            ss.supplier_db['magic_links'].append({"token": token, "type": "material", "target_id": mat_map.get(sel_mat_for_link, "")})
            st.success(f"Created token: {token}")
        links_df = pd.DataFrame(ss.supplier_db['magic_links']) if ss.supplier_db['magic_links'] else pd.DataFrame(columns=["token","type","target_id"]) 
        st.dataframe(links_df, use_container_width=True)

# -------------------------
# Assurance Readiness
# -------------------------
elif ss.selected_module == "Assurance Readiness":
    st.markdown("## ✅ Assurance Readiness Checklist")
    area = st.selectbox("Disclosure area", ["E2 – Pollution/SVHC","E3 – Water","E5 – Circular Economy"]) 
    cols = st.columns(3)
    with cols[0]:
        evidence = st.checkbox("Evidence attached (declarations/certificates)", True)
        trace = st.checkbox("Traceable to SKU/material", True)
        coverage = st.slider("Supplier/material coverage (%)", 0, 100, 70)
    with cols[1]:
        owner = st.text_input("Control owner", "Plant Quality Lead")
        cadence = st.selectbox("Control cadence", ["Monthly","Quarterly","Annually"]) 
        last_audit = st.date_input("Last internal audit")
    with cols[2]:
        soa = st.checkbox("Statement of applicability updated", False)
        sop = st.checkbox("SOPs current & trained", True)
        issues = st.number_input("Open audit issues", 0, 50, 1)
    score = ((20 if evidence else 0) + (20 if trace else 0) + (coverage * 0.3) + (10 if cadence in ["Monthly","Quarterly"] else 0) + (10 if sop else 0) - min(issues * 2, 10))
    score = max(0, min(100, round(score, 1)))
    band = "🟢 Ready" if score >= 80 else "🟡 Limited" if score >= 60 else "🔴 Not Ready"
    st.markdown(f"### Score: **{score}** / 100 — {band}")

# -------------------------
# Report Generator
# -------------------------
elif ss.selected_module == "Report Generator":
    st.markdown("### Board Presentation Draft")
    if Presentation is None:
        st.info("Install `python-pptx` to enable PPT export.")
    else:
        if st.button("📤 Download Board Deck (PPTX)"):
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "ESG Performance — Board Update (Demo)"
            slide.placeholders[1].text = "On Track: 92% diversion, €2.3M savings, −12% water YoY\nAttention: CSRD 78% (target 95%), €1.8M risk, Scope 3 gaps"
            s2 = prs.slides.add_slide(prs.slide_layouts[1])
            s2.shapes.title.text = "CSRD Readiness"
            tf = s2.placeholders[1].text_frame
            tf.text = "Current: 78% → Target: 95% by Q2 2025"
            tf.add_paragraph().text = "Risk exposure: €1.8M potential fines"
            tf.add_paragraph().text = "Mitigation: Supplier portal, SVHC tracking, assurance"
            s3 = prs.slides.add_slide(prs.slide_layouts[1])
            s3.shapes.title.text = "Waste Excellence"
            tf3 = s3.placeholders[1].text_frame
            tf3.text = "35% reduction possible; €2.3M annual savings"
            tf3.add_paragraph().text = "Deploy optimizer; replicate best practices"
            s4 = prs.slides.add_slide(prs.slide_layouts[1])
            s4.shapes.title.text = "Investment"
            tf4 = s4.placeholders[1].text_frame
            tf4.text = "Approve €750K to enable savings and compliance"
            bio = BytesIO(); prs.save(bio); bio.seek(0)
            st.download_button("📥 Download PPTX", data=bio.getvalue(), file_name="Board_Update_Demo.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

st.markdown("---")
st.caption("P&G ESG Intelligence Platform — AI Kit A Demo • This is simulated demo data for presentation purposes.")
